import os
import io
import json
import base64
import tempfile
import asyncio
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse

from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi import status
from pptx import Presentation
import httpx
from dotenv import load_dotenv
from openai import AsyncOpenAI
from typing import TypedDict,Optional
import os 

load_dotenv()

class State(TypedDict):
    mode : str
    file_path : Optional[str]
    github_url : Optional[str]
    video_url : Optional[str]
    content: str
    output: Optional[dict] | None

load_dotenv()



OPEN_AI_KEY= os.getenv("OPENAI_API_KEY", "")

OPENAI_API_KEY= OPEN_AI_KEY
if not OPENAI_API_KEY:
    raise RuntimeError("OPEN_AI_KEY is not set")

client = AsyncOpenAI(api_key=OPENAI_API_KEY)

router = APIRouter( tags=["ppt"])

GPT_MODEL = "gpt-4o-mini"
HTTP_TIMEOUT = 30
MAX_CONCURRENT_SLIDES = 4


def img_to_b64(blob: bytes) -> str:
    return "data:image/png;base64," + base64.b64encode(blob).decode("utf-8")


async def load_presentation(source: str) -> Presentation:
    if source.startswith("http://") or source.startswith("https://"):
        async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as c:
            r = await c.get(source)
            if r.status_code != 200:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Failed to fetch PPT from URL, status {r.status_code}",
                )
            return Presentation(io.BytesIO(r.content))
    return Presentation(source)


async def extract_ppt_slides(source: str) -> List[Dict[str, Any]]:
    prs = await load_presentation(source)
    slides: List[Dict[str, Any]] = []
    for idx, s in enumerate(prs.slides):
        texts: List[str] = []
        images: List[str] = []
        for sh in s.shapes:
            if hasattr(sh, "text") and isinstance(sh.text, str):
                t = sh.text.strip()
                if t:
                    texts.append(t)
            if getattr(sh, "shape_type", None) == 13 and hasattr(sh, "image"):
                try:
                    images.append(img_to_b64(sh.image.blob))
                except Exception:
                    continue
        slides.append(
            {
                "index": idx + 1,
                "text": "\n".join(texts),
                "images": images,
            }
        )
    return slides


async def call_gpt_json(messages: List[Dict[str, Any]]) -> Dict[str, Any]:
    try:
        res = await client.chat.completions.create(
            model=GPT_MODEL,
            messages=messages,
            response_format={"type": "json_object"},
            temperature=0,
            timeout=HTTP_TIMEOUT,
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"LLM call failed: {str(e)}",
        )
    content = res.choices[0].message.content
    try:
        return json.loads(content)
    except Exception:
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail="LLM returned invalid JSON in JSON mode",
        )


async def analyze_single_slide(topic: str, slide: Dict[str, Any]) -> Dict[str, Any]:
    rubric = (
        "You are a professional pitch-deck evaluator. "
        "Return ONLY valid JSON in the EXACT structure below, no explanation, no prose:\n\n"
        "{\n"
        ' "clarity": {\n'
        '   "headline_present": bool,\n'
        '   "key_message_present": bool,\n'
        '   "text_density": "low" | "medium" | "high",\n'
        '   "readability_score": number\n'
        " },\n"
        ' "design": {\n'
        '   "alignment_good": bool,\n'
        '   "contrast_good": bool,\n'
        '   "visual_hierarchy": "strong" | "weak",\n'
        '   "consistency_issues": [string]\n'
        " },\n"
        ' "storytelling": {\n'
        '   "problem_defined": bool,\n'
        '   "solution_defined": bool,\n'
        '   "use_case_clear": bool,\n'
        '   "logical_flow": "yes" | "no"\n'
        " },\n"
        ' "missing_elements": [\n'
        '    "architecture", "market_analysis", "competitors",\n'
        '    "demo", "business_model", "persona", "roadmap"\n'
        " ],\n"
        ' "issues_detected": [string],\n'
        ' "manipulation_detected": bool,\n'
        ' "suggestions": [string]\n'
        "}\n"
        "Only include items in missing_elements that are truly missing for this slide."
    )

    content_parts: List[Dict[str, Any]] = [
        {
            "type": "text",
            "text": f"Topic: {topic}\nSlide {slide['index']}:\n{slide['text']}",
        }
    ]

    for img in slide.get("images", [])[:4]:
        content_parts.append(
            {
                "type": "image_url",
                "image_url": {"url": img},
            }
        )

    messages = [
        {"role": "system", "content": rubric},
        {"role": "user", "content": content_parts},
    ]

    return await call_gpt_json(messages)


def compute_slide_scores(analysis: Dict[str, Any]) -> Dict[str, float]:
    clarity = analysis.get("clarity") or {}
    design = analysis.get("design") or {}
    storytelling = analysis.get("storytelling") or {}

    clarity_raw = 0.0
    clarity_max = 4.0
    if clarity.get("headline_present") is True:
        clarity_raw += 1
    if clarity.get("key_message_present") is True:
        clarity_raw += 1
    td = clarity.get("text_density")
    if td == "low":
        clarity_raw += 1
    elif td == "medium":
        clarity_raw += 0.5
    rs = clarity.get("readability_score")
    if isinstance(rs, (int, float)):
        r_norm = max(0.0, min(float(rs), 100.0)) / 100.0
        clarity_raw += r_norm
    clarity_score = 0.0 if clarity_max == 0 else (clarity_raw / clarity_max) * 100.0

    design_raw = 0.0
    design_max = 3.0
    if design.get("alignment_good") is True:
        design_raw += 1
    if design.get("contrast_good") is True:
        design_raw += 1
    if design.get("visual_hierarchy") == "strong":
        design_raw += 1
    design_score = 0.0 if design_max == 0 else (design_raw / design_max) * 100.0

    story_raw = 0.0
    story_max = 4.0
    if storytelling.get("problem_defined") is True:
        story_raw += 1
    if storytelling.get("solution_defined") is True:
        story_raw += 1
    if storytelling.get("use_case_clear") is True:
        story_raw += 1
    if storytelling.get("logical_flow") == "yes":
        story_raw += 1
    story_score = 0.0 if story_max == 0 else (story_raw / story_max) * 100.0

    available = []
    if clarity:
        available.append(clarity_score)
    if design:
        available.append(design_score)
    if storytelling:
        available.append(story_score)
    overall = sum(available) / len(available) if available else 0.0

    return {
        "clarity": round(clarity_score, 2),
        "design": round(design_score, 2),
        "story": round(story_score, 2),
        "overall": round(overall, 2),
    }


async def deck_level_analysis(topic: str, slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
    rubric = (
        "You evaluate the FULL PPT DECK. "
        "Return ONLY strict JSON in the structure below, no explanation:\n\n"
        "{\n"
        ' "missing_critical_sections": [string],\n'
        ' "strengths": [string],\n'
        ' "weaknesses": [string],\n'
        ' "narrative_flow": "good" | "moderate" | "broken",\n'
        ' "story_completeness": "strong" | "weak" | "incomplete",\n'
        ' "recommended_fixes": [string]\n'
        "}\n"
    )

    text_joined = "\n\n".join(
        [f"Slide {s['index']}:\n{s['text']}" for s in slides_data]
    )

    messages = [
        {"role": "system", "content": rubric},
        {
            "role": "user",
            "content": f"TOPIC: {topic}\n\nFULL_DECK:\n{text_joined}",
        },
    ]

    return await call_gpt_json(messages)


async def analyze_ppt_with_gpt(state: State) -> State:
    topic = state["content"]
    file_path = state["file_path"]

    slides = await extract_ppt_slides(file_path)

    sem = asyncio.Semaphore(MAX_CONCURRENT_SLIDES)
    slide_results: List[Dict[str, Any]] = []

    async def process_slide(slide: Dict[str, Any]) -> None:
        async with sem:
            try:
                analysis = await analyze_single_slide(topic, slide)
                scores = compute_slide_scores(analysis)
                slide_results.append(
                    {
                        "slide_number": slide["index"],
                        "analysis": analysis,
                        "score": scores["overall"],
                        "score_breakdown": scores,
                    }
                )
            except HTTPException:
                raise
            except Exception as e:
                slide_results.append(
                    {
                        "slide_number": slide["index"],
                        "analysis": {"error": str(e)},
                        "score": 0.0,
                        "score_breakdown": {
                            "clarity": 0.0,
                            "design": 0.0,
                            "story": 0.0,
                            "overall": 0.0,
                        },
                    }
                )

    await asyncio.gather(*(process_slide(slide) for slide in slides))

    slide_results_sorted = sorted(slide_results, key=lambda x: x["slide_number"])
    mentor_summary = await generate_human_readable_mentorship(topic, slide_results_sorted)
    valid_scores = [s for s in slide_results_sorted if s["score"] is not None]
    if valid_scores:
        clarity_avg = sum(s["score_breakdown"]["clarity"] for s in valid_scores) / len(
            valid_scores
        )
        design_avg = sum(s["score_breakdown"]["design"] for s in valid_scores) / len(
            valid_scores
        )
        story_avg = sum(s["score_breakdown"]["story"] for s in valid_scores) / len(
            valid_scores
        )
    else:
        clarity_avg = design_avg = story_avg = 0.0

    deck_summary = await deck_level_analysis(topic, slides)

    missing_critical = deck_summary.get("missing_critical_sections") or []
    completeness_score = max(0.0, 100.0 - len(missing_critical) * 10.0)

    overall = round(
        (clarity_avg + design_avg + story_avg + completeness_score) / 4.0, 2
    )

    final_feedback = (
        f"The deck has strengths in {', '.join(deck_summary.get('strengths', []))}. "
        f"However, weaknesses include {', '.join(deck_summary.get('weaknesses', []))}. "
        f"Critical missing sections: {', '.join(missing_critical)}. "
        f"Recommended fixes: {', '.join(deck_summary.get('recommended_fixes', []))}."
    )

    output = {
        "slides": slide_results_sorted,
        "deck_summary": deck_summary,
        "score": {
            "clarity_score": round(clarity_avg, 2),
            "design_score": round(design_avg, 2),
            "story_score": round(story_avg, 2),
            "completeness_score": round(completeness_score, 2),
            "overall_score": overall,
        },
        "final_feedback": final_feedback,
        "mentor_summary": mentor_summary,
    }

    return {
        **state,
        "output": output,
    }

async def generate_human_readable_mentorship(topic: str, slides: List[Dict[str, Any]]) -> str:
    slides_text = ""
    for s in slides:
        slides_text += f"Slide {s['slide_number']}:\n{s['analysis']}\n\n"

    system_prompt = (
        "You are a world-class PPT mentor and must output ONLY MARKDOWN. "
        "Use EXACTLY this format:\n\n"
        "# <PPT Title> – Analysis\n\n"
        "## Slide 1\n"
        "###  Problem\n"
        "<problem>\n\n"
        "###  How to Fix\n"
        "<fix>\n\n"
        "---\n\n"
        "(repeat for every slide)\n\n"
        "## Conclusion\n"
        "<overall summary>\n\n"
        "RULES:\n"
        "- MUST use Markdown headings (#, ##, ###).\n"
        "- MUST use  and  emojis.\n"
        "- MUST insert '---' after each slide.\n"
        "- No extra sections.\n"
        "- No intro, no explanation.\n"
        "- Provide clear, human-friendly actionable points.\n"
    )

    user_prompt = f"Topic: {topic}\n\nSlide Analysis:\n{slides_text}"

    res = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0
    )

    return res.choices[0].message.content



@router.post("/analyze")
async def analyze_ppt_endpoint(
    topic: str = Form(...),
    file: Optional[UploadFile] = File(None),
    file_url: Optional[str] = Form(None),
): 
    print("Received analyze request with topic:", topic)
    if not file and not file_url:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Either file or file_url is required",
        )

    temp_path: Optional[str] = None

    if file:
        filename = file.filename or ""
        ext = os.path.splitext(filename)[1].lower()
        if ext not in [".pptx", ".ppt"]:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Only .ppt or .pptx files are supported",
            )
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            content = await file.read()
            tmp.write(content)
            temp_path = tmp.name
        file_path = temp_path
    else:
        parsed = urlparse(file_url)
        if parsed.scheme not in ("http", "https"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="file_url must be http or https",
            )
        file_path = file_url

    state: State = {
        "content": topic,
        "file_path": file_path,
    }

    try:
        result_state = await analyze_ppt_with_gpt(state)
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass

    return result_state["output"]

@router.get("/health")
async def health_check():
    return {"status": "ok"}
